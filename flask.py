import requests
import os
import re
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from pptx import Presentation
from jinja2 import Template

app = Flask(__name__)

# クロスオリジンリソースシェアリングの設定
CORS(app, resources={r"/*": {"origins": r"https://.*\.xxxx\.com"}})

# 環境変数から設定値を取得
CERTIFICATE_FILE_PATH = os.getenv('FLASK_SERVER_CERTIFICATE_FILE_PATH', 'path/to/default/cert.pem')
KEY_FILE_PATH = os.getenv('FLASK_SERVER_KEY_FILE_PATH', 'path/to/default/key.pem')
OUTPUT_FILE_PATH = os.getenv('OUTPUT_FILE_PATH', 'path/to/default/output_presentation.pptx')
AI_GATEWAY_ENDPOINT = os.environ.get('AI_GATEWAY_ENDPOINT')

# AI Gatewayへのリクエストのヘッダー
headers = {
    "Authorization": "Bearer xxxxxxx", 
    "Content-Type": "application/json"
}

# プロンプト
prompt_outline_generate = """
You are a power point presentation specialist. You are asked to create
the outline for a presentation in japanese about {{ topic }}.

Try your best to create a outline that fits the needs as follows.
1. logical
2. understandable
3. fitable for a presentation
Return the outline as a Markdown format and mark each part with the comment of <!-- title -->, <!-- head section --> or <!-- title and content -->.

Markdown Examples:
<!-- title -->
# 運動しましょう

<!-- head section -->
## 運動の種類

<!-- title and content -->
### チームでプレーできる運動
バレーボール
バスケットボール

<!-- head section -->
## 運動の重要性

<!-- title and content -->
### 身体的健康への利点
慢性疾患のリスク低減
筋力、柔軟性、持久力の向上

<!-- title and content -->
### 精神的健康への利点
精神疾患の症状軽減
気分の向上
認知機能の改善

<!-- head section -->
## 終わり
"""
prompt_layout_classfication = """
You are a power point presentation specialist.

Please extract 3 layouts with 'title', 'head section', 'title and content' layout from layout information as follows.
---
{{ layouts_info }}
---

return a markdown format like the example as follows.
you have include comment with 'title', 'head section', or 'title and content' , Layout (Index,Name) and Placeholder (Index, Name, Type, Area) as follows.

Example:
<!-- title -->
Layout Index: 0, Name: Title
- Placeholder Index: 0, Name: Title 1, Type: TITLE (1), Area: 100 square centimeters　
- Placeholder Index: 1, Name: Picture Placeholder 2, Type: PICTURE (18), Area: 200 square centimeters　

<!-- head section -->
Layout Index: 3, Name:
- Placeholder Index: 0, Name: Title 1, Type: TITLE (1), Area: 50 square centimeters
- Placeholder Index: 1, Name: Picture Placeholder 2, Type: PICTURE (18), Area: 500 square centimeters

<!-- title and content -->
Layout Index: 3, Name:
<!-- body -->
- Placeholder Index: 0, Name: Title 1, Type: TITLE (1), Area: 223 square centimeters
- Placeholder Index: 1, Name: Picture Placeholder 2, Type: PICTURE (18), Area: 500 square centimeters
"""
prompt_slides_json_generate = """
You are a power point presentation specialist. You are asked to create
the content for a presentation in japanese based on the outline provided.
outline:
---
{{ outline_generated }}
---

Please create a JSON including the slides' content using the outline and stick to the layoutinformation as follows.

layout information
---
{{ layout_category }}
---
Return the structured JSON like the example as follows.
Your answer should contain layout_index, placeholder index,and the information input in the placeholder.

example:
{{ slides_json_example }}
"""
slides_json_example = """
{
    "slides": [
        {
            "layout_index": 0,
            "content": [
                {"placeholder_index": 0, "placeholder_text": "関数"},
            ]
        },
        {
            "layout_index": 1,
            "content": [
                {"placeholder_index": 0, "placeholder_text": "関数とは"}
            ]
        },
        {
            "layout_index": 3,
            "content": [
                {"placeholder_index": 0, "placeholder_text": "関数の定義"},
                {"placeholder_index": 14, "placeholder_text": "よく行う処理をひとまとめにしたもの\n関数を呼び出すことで、まとめた処理が実行される"}
            ]
        },

    ]
}
"""
prompt_code_generate = """
"You are a PowerPoint presentation specialist. You will get a list of slides by JSON. "
"You need to create a PowerPoint presentation called \"output_presentation.pptx\" based on the provided slides. "
"But there is a catch: Instead of creating the presentation, provide Python code that generates the PowerPoint presentation. Use the package python-pptx."
"In your code, use a file called \"template.pptx\" as the template for the presentation and stick to the template's design.\n\n"

"For the slide contents. Use the information of layout_index, placeholder_index, and placeholder_text in the Slides as follows. "

"Slides:{{ outline_json }}"

"Your answer should only contain the Python code, no explanatory text.\n\n"

"""

# 関数
def generate_outline(topic):
    """
    Generate an outline based on the given topic using GPT.
    """
    print(f"Starting generate_outline with topic: {topic}") 

    prompt = Template(prompt_outline_generate).render(topic=topic)
    print(f"Generated prompt: {prompt}")  

    payload = {
        "model": "azure/gpt4o-mini",
        "messages": [
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
        ]
    }
    
    print(f"Payload to be sent: {payload}")  

    try:
        response = requests.post(f"{AI_GATEWAY_ENDPOINT}/chat/completions", headers=headers, json=payload)
        print(f"Response Status Code: {response.status_code}")  

        if response.status_code != 200:
            print(f"Error: Received status code {response.status_code} from the server.")
            print(f"Response content: {response.text}")
            return None

        response_data = response.json()
        print(f"Response JSON: {response_data}") 

        if 'choices' in response_data and len(response_data['choices']) > 0:
            outline = response_data['choices'][0]['message']['content']
            return outline
        else:
            print("No valid outline found in the response.")
            return None

    except requests.exceptions.RequestException as e:
        print(f"Request error: {e}")
        return None
def request_code_and_execute(payload, attempt=1):
    """
    Sends the payload to GPT, extracts the generated Python code, and executes it.
    If an error occurs, it retries once more by sending the error back to GPT.
    """
    response = requests.post(f"{AI_GATEWAY_ENDPOINT}/chat/completions", headers=headers, json=payload)
    presentation_code = response.json()['choices'][0]['message']['content']
    match = re.findall(r"```python\n(.*?)\n```", presentation_code, re.DOTALL)

    if not match:
        print("No valid Python code found.")
        return

    python_code = match[0]

    try:
        exec(python_code)
        print("Presentation generated successfully.")
        return python_code
    except Exception as e:
        print(f"Error occurred during execution: {e}")

        if attempt == 1:
            error_message = f"An error occurred while executing the generated code: {e}"
            new_payload = {
                "model": "azure/gpt4o-mini",
                "messages": [
                    {"role": "system", "content": "You are a helpful assistant."},
                    {
                        "role": "user",
                        "content": (
                            f"The code you provided caused the following error:\n\n{e}\n\n"
                            "Please revise the Python code to address the error. The updated code should "
                            "still follow the original instructions, including loading the template 'template.pptx', "
                            "generating slides based on the provided slide data."
                        )
                    }
                ]
            }

            request_code_and_execute(new_payload, attempt=2)
        else:
            print("The second attempt also failed. Here is the code generated by GPT:")
            print(python_code)
def emu_to_cm(emu):
    return emu / 360000 * 2.54
def calculate_area(width, height):
    width_cm = emu_to_cm(width)
    height_cm = emu_to_cm(height)
    return width_cm * height_cm
def layout_categorize(layouts_info):
    prompt = Template(prompt_layout_classfication).render(layouts_info=layouts_info)
    print(f"Generated prompt: {prompt}")  

    payload = {
        "model": "azure/gpt4o-mini",
        "messages": [
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
        ]
    }
    
    print(f"Payload to be sent: {payload}")  

    try:
        response = requests.post(f"{AI_GATEWAY_ENDPOINT}/chat/completions", headers=headers, json=payload)
        print(f"Response Status Code: {response.status_code}") 

        if response.status_code != 200:
            print(f"Error: Received status code {response.status_code} from the server.")
            print(f"Response content: {response.text}")
            return None

        response_data = response.json()
        print(f"Response JSON: {response_data}")  

        if 'choices' in response_data and len(response_data['choices']) > 0:
            layout_category = response_data['choices'][0]['message']['content']
            return layout_category
        else:
            print("No valid outline found in the response.")
            return None

    except requests.exceptions.RequestException as e:
        print(f"Request error: {e}")
        return None
    

# プレゼンテーション生成のエンドポイント
@app.route('/pptx_generation', methods=['POST'])
def layout_output():
    if 'template_file' not in request.files:
        return jsonify({"error": "No file part"}), 400

    template_file = request.files['template_file']
    if template_file:
        template_file_path = os.path.join("./", "template.pptx")
        template_file.save(template_file_path)

    prs = Presentation('template.pptx')

    if template_file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    topic = request.form.get('topic')
    if not topic:
        return jsonify({"error": "No topic provided"}), 400

    outline_generated = generate_outline(topic)
    print(f"outline_generated: {outline_generated}")
    if not outline_generated:
        print("Failed to generate outline.")
        return jsonify({"error": "Failed to generate outline"}), 500

    layout_infos = []

    for layout_index, layout in enumerate(prs.slide_master.slide_layouts):
        layout_info = {
            "layout_index": layout_index,
            "layout_name": layout.name,
            "placeholders": []
        }

        for shape in layout.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                ph_idx = shape.placeholder_format.idx
                ph_name = shape.name
                ph_width = shape.width
                ph_height = shape.height
                area = calculate_area(ph_width, ph_height)

                layout_info["placeholders"].append({
                    "index": ph_idx,
                    "name": ph_name,
                    "type": ph_type.name,
                    "type_value": ph_type.value,
                    "area": area
                })

        layout_infos.append(layout_info)

    layout_category = layout_categorize(layout_infos)

    prompt = Template(prompt_slides_json_generate).render(
        layout_category=layout_category,
        outline_generated=outline_generated,
        slides_json_example=slides_json_example
    )

    payload = {
        "model": "azure/gpt4o-mini",
        "messages": [
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
        ]
    }

    try:
        response = requests.post(f"{AI_GATEWAY_ENDPOINT}/chat/completions", headers=headers, json=payload)
        print(f"Response Status Code: {response.status_code}")

        if response.status_code != 200:
            print(f"Error: Received status code {response.status_code} from the server.")
            print(f"Response content: {response.text}")
            return jsonify({"error": "Failed to generate JSON from AI response."}), 500

        response_data = response.json()

        if 'choices' in response_data and len(response_data['choices']) > 0:
            outline_json = response_data['choices'][0]['message']['content']
            print(f"Outline JSON: {outline_json}")

            prompt_pptx = Template(prompt_code_generate).render(
                outline_json=outline_json
            )

            payload_pptx_generate = {
                "model": "azure/gpt4o-mini",
                "messages": [
                    {"role": "system", "content": "You are a helpful assistant."},
                    {"role": "user", "content": prompt_pptx}
                ]
            }
            
            resp = request_code_and_execute(payload_pptx_generate)
            print(f"Code_11111: {resp}")

            if not os.path.exists(OUTPUT_FILE_PATH):
                print(f"Error: File not found at {OUTPUT_FILE_PATH}. Please check the file generation process.")
                return jsonify({"error": "Failed to generate the presentation file."}), 500
            return send_file(
                OUTPUT_FILE_PATH,
                as_attachment=True,
                download_name="output_file_sample.pptx",
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )

        else:
            print("No valid outline found in the response.")
            return jsonify({"error": "No valid outline found in the AI response."}), 500

    except requests.exceptions.RequestException as e:
        print(f"Request error: {e}")
        return jsonify({"error": "Request to AI Gateway failed."}), 500


if __name__ == '__main__':
    # SSL設定を含むサーバーの実行
    app.run(debug=True, ssl_context=(CERTIFICATE_FILE_PATH, KEY_FILE_PATH), host='127.0.0.1', port=5000)
